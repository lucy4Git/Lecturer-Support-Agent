from __future__ import annotations

import csv
import hashlib
import json
import re
from io import BytesIO, StringIO
from pathlib import Path
from typing import Callable

from bs4 import BeautifulSoup
from charset_normalizer import from_bytes
from docx import Document as WordDocument
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

from .contracts import ExtractionResult


class DocumentParserRegistry:
    def __init__(self) -> None:
        self._suffix_parsers: dict[str, Callable[[bytes], ExtractionResult]] = {
            ".txt": self._plain_text,
            ".md": self._plain_text,
            ".markdown": self._plain_text,
            ".csv": self._csv,
            ".json": self._json,
            ".html": self._html,
            ".htm": self._html,
            ".pdf": self._pdf,
            ".docx": self._docx,
            ".pptx": self._pptx,
            ".xlsx": self._xlsx,
            ".vtt": self._transcript,
            ".srt": self._transcript,
        }

    def extract(self, *, filename: str, media_type: str | None, content: bytes) -> ExtractionResult:
        suffix = Path(filename).suffix.lower()
        if suffix in {".mp3", ".wav", ".m4a", ".mp4", ".mov", ".webm", ".avi"} or (
            media_type and (media_type.startswith("audio/") or media_type.startswith("video/"))
        ):
            return ExtractionResult(
                status="transcript_required",
                parser_name="media_metadata_v1",
                transcript_required=True,
                warnings=["Audio and video files require an authorised transcript before semantic indexing."],
                metadata={"media_type": media_type, "content_sha256": hashlib.sha256(content).hexdigest()},
            )
        if suffix in {".png", ".jpg", ".jpeg", ".gif", ".webp", ".tif", ".tiff"} or (
            media_type and media_type.startswith("image/")
        ):
            return ExtractionResult(
                status="unsupported",
                parser_name="image_metadata_v1",
                warnings=["Image OCR is not enabled in v1.6; the original image is preserved without invented text."],
                metadata={"media_type": media_type, "content_sha256": hashlib.sha256(content).hexdigest()},
            )
        parser = self._suffix_parsers.get(suffix)
        if parser is None:
            return ExtractionResult(
                status="unsupported",
                parser_name="unsupported_format_v1",
                warnings=[f"No safe text extractor is configured for {suffix or media_type or 'this format'}."],
            )
        try:
            result = parser(content)
        except Exception as exc:
            return ExtractionResult(
                status="failed",
                parser_name=f"{suffix.lstrip('.') or 'unknown'}_parser_v1",
                warnings=[f"Extraction failed: {type(exc).__name__}."],
                metadata={"safe_error_type": type(exc).__name__},
            )
        clean = self._normalise(result.text)
        status = result.status if clean else ("empty" if result.status == "extracted" else result.status)
        return result.model_copy(
            update={
                "text": clean,
                "status": status,
                "word_count": len(re.findall(r"\b\w+\b", clean)),
                "quality_score": self._quality(clean),
            }
        )

    @staticmethod
    def _decode(content: bytes) -> str:
        match = from_bytes(content).best()
        if match is None:
            return content.decode("utf-8", errors="replace")
        return str(match)

    @staticmethod
    def _normalise(text: str) -> str:
        text = text.replace("\x00", "")
        text = re.sub(r"\r\n?", "\n", text)
        text = re.sub(r"[ \t]+\n", "\n", text)
        text = re.sub(r"\n{4,}", "\n\n\n", text)
        return text.strip()

    @staticmethod
    def _quality(text: str) -> str:
        if len(text) >= 1000:
            return "high"
        if len(text) >= 100:
            return "medium"
        return "low"

    def _plain_text(self, content: bytes) -> ExtractionResult:
        return ExtractionResult(status="extracted", parser_name="plain_text_v1", text=self._decode(content))

    def _transcript(self, content: bytes) -> ExtractionResult:
        text = self._decode(content)
        text = re.sub(r"^WEBVTT.*?$", "", text, flags=re.MULTILINE)
        text = re.sub(r"^\d+$", "", text, flags=re.MULTILINE)
        text = re.sub(r"\d{1,2}:\d{2}:\d{2}[,\.]\d{3}\s+-->\s+.*$", "", text, flags=re.MULTILINE)
        return ExtractionResult(status="extracted", parser_name="transcript_v1", text=text)

    def _csv(self, content: bytes) -> ExtractionResult:
        decoded = self._decode(content)
        rows = list(csv.reader(StringIO(decoded)))
        text = "\n".join(" | ".join(cell.strip() for cell in row) for row in rows)
        return ExtractionResult(status="extracted", parser_name="csv_v1", text=text, table_count=1)

    def _json(self, content: bytes) -> ExtractionResult:
        value = json.loads(self._decode(content))
        return ExtractionResult(
            status="extracted", parser_name="json_v1", text=json.dumps(value, ensure_ascii=False, indent=2)
        )

    def _html(self, content: bytes) -> ExtractionResult:
        soup = BeautifulSoup(self._decode(content), "html.parser")
        for node in soup(["script", "style", "noscript"]):
            node.decompose()
        return ExtractionResult(status="extracted", parser_name="html_v1", text=soup.get_text("\n"))

    @staticmethod
    def _pdf(content: bytes) -> ExtractionResult:
        reader = PdfReader(BytesIO(content))
        pages = [(page.extract_text() or "").strip() for page in reader.pages]
        text = "\n\n".join(f"[Page {index}]\n{page}" for index, page in enumerate(pages, start=1) if page)
        warnings = [] if text else ["The PDF contained no extractable text; it may require OCR."]
        return ExtractionResult(
            status="extracted", parser_name="pypdf_v1", text=text, page_count=len(reader.pages), warnings=warnings
        )

    @staticmethod
    def _docx(content: bytes) -> ExtractionResult:
        document = WordDocument(BytesIO(content))
        parts: list[str] = []
        for paragraph in document.paragraphs:
            if paragraph.text.strip():
                parts.append(paragraph.text.strip())
        table_count = 0
        for table in document.tables:
            table_count += 1
            for row in table.rows:
                parts.append(" | ".join(cell.text.strip() for cell in row.cells))
        return ExtractionResult(
            status="extracted", parser_name="python_docx_v1", text="\n\n".join(parts), table_count=table_count
        )

    @staticmethod
    def _pptx(content: bytes) -> ExtractionResult:
        presentation = Presentation(BytesIO(content))
        slides: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            text = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            if text:
                slides.append(f"[Slide {index}]\n" + "\n".join(text))
        return ExtractionResult(
            status="extracted", parser_name="python_pptx_v1", text="\n\n".join(slides), slide_count=len(presentation.slides)
        )

    @staticmethod
    def _xlsx(content: bytes) -> ExtractionResult:
        workbook = load_workbook(BytesIO(content), read_only=True, data_only=True)
        output: list[str] = []
        table_count = 0
        for sheet in workbook.worksheets:
            output.append(f"[Sheet: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                values = ["" if value is None else str(value) for value in row]
                if any(values):
                    output.append(" | ".join(values))
            table_count += 1
        return ExtractionResult(
            status="extracted",
            parser_name="openpyxl_v1",
            text="\n".join(output),
            sheet_count=len(workbook.worksheets),
            table_count=table_count,
        )
