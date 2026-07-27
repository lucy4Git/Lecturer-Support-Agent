from __future__ import annotations

import html
import re
from copy import copy
from dataclasses import dataclass
from io import BytesIO

from docx import Document as WordDocument
from openpyxl import Workbook
from pptx import Presentation
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

from services.database.models.enums import ExportAudience, ExportFormat


@dataclass(frozen=True, slots=True)
class RenderedExport:
    content: bytes
    media_type: str
    extension: str


class ExportRenderer:
    """Render a current immutable output version into common teaching formats."""

    _answer_heading = re.compile(
        r"(?i)^#{1,6}\s*(answer key|answers|marking guide|memorandum|model answers?|solutions?|"
        r"lecturer notes?|marking notes?|moderator notes?|assessment rationale)\s*$"
    )
    _inline_confidential = re.compile(
        r"(?i)^\s*(answer|solution|expected answer|model answer|marking note|lecturer note)\s*:\s*"
    )

    def prepare_content(self, markdown: str, audience: ExportAudience) -> str:
        if audience != ExportAudience.STUDENT_COPY:
            return markdown
        lines = markdown.splitlines()
        kept: list[str] = []
        suppress = False
        for line in lines:
            stripped = line.strip()
            if self._answer_heading.match(stripped):
                suppress = True
                continue
            if suppress and re.match(r"^#{1,6}\s+", line):
                suppress = False
            if suppress:
                continue
            if self._inline_confidential.match(stripped):
                continue
            # Remove compact bracketed answer disclosures without altering the question text.
            line = re.sub(
                r"(?i)\s*\[(?:answer|solution|model answer)\s*:\s*[^\]]+\]\s*$",
                "",
                line,
            )
            kept.append(line)
        return "\n".join(kept).strip()

    def render(
        self,
        *,
        title: str,
        markdown: str,
        export_format: ExportFormat,
        audience: ExportAudience,
    ) -> RenderedExport:
        content = self.prepare_content(markdown, audience)
        if export_format == ExportFormat.MARKDOWN:
            return RenderedExport(content.encode("utf-8"), "text/markdown; charset=utf-8", "md")
        if export_format == ExportFormat.HTML:
            return RenderedExport(self._html(title, content), "text/html; charset=utf-8", "html")
        if export_format == ExportFormat.DOCX:
            return RenderedExport(self._docx(title, content), "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "docx")
        if export_format == ExportFormat.PDF:
            return RenderedExport(self._pdf(title, content), "application/pdf", "pdf")
        if export_format == ExportFormat.PPTX:
            return RenderedExport(self._pptx(title, content), "application/vnd.openxmlformats-officedocument.presentationml.presentation", "pptx")
        if export_format == ExportFormat.XLSX:
            return RenderedExport(self._xlsx(title, content), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "xlsx")
        raise ValueError(f"Unsupported export format: {export_format}")

    def _html(self, title: str, markdown: str) -> bytes:
        rows: list[str] = []
        for line in markdown.splitlines():
            escaped = html.escape(line)
            if line.startswith("### "):
                rows.append(f"<h3>{html.escape(line[4:])}</h3>")
            elif line.startswith("## "):
                rows.append(f"<h2>{html.escape(line[3:])}</h2>")
            elif line.startswith("# "):
                rows.append(f"<h1>{html.escape(line[2:])}</h1>")
            elif line.startswith(("- ", "* ")):
                rows.append(f"<p>• {html.escape(line[2:])}</p>")
            elif line.strip():
                rows.append(f"<p>{escaped}</p>")
        document = (
            "<!doctype html><html><head><meta charset='utf-8'><title>"
            + html.escape(title)
            + "</title></head><body>"
            + "".join(rows)
            + "</body></html>"
        )
        return document.encode("utf-8")

    def _docx(self, title: str, markdown: str) -> bytes:
        document = WordDocument()
        document.add_heading(title, level=0)
        for line in markdown.splitlines():
            if line.startswith("### "):
                document.add_heading(line[4:], level=3)
            elif line.startswith("## "):
                document.add_heading(line[3:], level=2)
            elif line.startswith("# "):
                document.add_heading(line[2:], level=1)
            elif line.startswith(("- ", "* ")):
                document.add_paragraph(line[2:], style="List Bullet")
            elif re.match(r"^\d+\. ", line):
                document.add_paragraph(re.sub(r"^\d+\. ", "", line), style="List Number")
            elif line.strip():
                document.add_paragraph(line)
        buffer = BytesIO(); document.save(buffer); return buffer.getvalue()

    def _pdf(self, title: str, markdown: str) -> bytes:
        buffer = BytesIO()
        styles = getSampleStyleSheet()
        story = [Paragraph(html.escape(title), styles["Title"]), Spacer(1, 5 * mm)]
        for line in markdown.splitlines():
            if not line.strip():
                story.append(Spacer(1, 2 * mm)); continue
            if line.startswith("### "):
                style = styles["Heading3"]; text = line[4:]
            elif line.startswith("## "):
                style = styles["Heading2"]; text = line[3:]
            elif line.startswith("# "):
                style = styles["Heading1"]; text = line[2:]
            elif line.startswith(("- ", "* ")):
                style = styles["BodyText"]; text = "• " + line[2:]
            else:
                style = styles["BodyText"]; text = line
            story.append(Paragraph(html.escape(text), style))
            story.append(Spacer(1, 1.5 * mm))
        SimpleDocTemplate(buffer, pagesize=A4, rightMargin=18*mm, leftMargin=18*mm, topMargin=18*mm, bottomMargin=18*mm).build(story)
        return buffer.getvalue()

    def _pptx(self, title: str, markdown: str) -> bytes:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = "Generated by Lecturer Support Agent — authorised human review remains required where applicable."
        sections = self._sections(markdown)
        for section_title, body in sections:
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = section_title[:120]
            frame = slide.placeholders[1].text_frame
            frame.clear()
            for index, line in enumerate([item for item in body.splitlines() if item.strip()] or [""]):
                paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
                paragraph.text = re.sub(r"^[-*]\s+", "", line)[:1200]
        buffer=BytesIO(); presentation.save(buffer); return buffer.getvalue()

    def _xlsx(self, title: str, markdown: str) -> bytes:
        workbook = Workbook(); sheet = workbook.active; sheet.title = "Teaching Output"
        sheet.append(["Title", title]); sheet.append([]); sheet.append(["Section", "Content"])
        for section_title, body in self._sections(markdown):
            sheet.append([section_title, body])
        sheet.column_dimensions["A"].width = 34; sheet.column_dimensions["B"].width = 110
        for row in sheet.iter_rows():
            for cell in row: alignment = copy(cell.alignment); alignment.wrap_text = True; alignment.vertical = "top"; cell.alignment = alignment
        buffer=BytesIO(); workbook.save(buffer); return buffer.getvalue()

    @staticmethod
    def _sections(markdown: str) -> list[tuple[str, str]]:
        sections: list[tuple[str, list[str]]] = [("Overview", [])]
        for line in markdown.splitlines():
            match = re.match(r"^#{1,6}\s+(.+)$", line)
            if match:
                sections.append((match.group(1).strip(), []))
            else:
                sections[-1][1].append(line)
        return [(title, "\n".join(lines).strip()) for title, lines in sections if lines or title != "Overview"]


def safe_export_filename(title: str, extension: str) -> str:
    stem = re.sub(r"[^A-Za-z0-9._-]+", "_", title).strip("._")[:120] or "teaching_output"
    return f"{stem}.{extension}"
